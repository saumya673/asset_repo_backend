from pptx import Presentation
from io import BytesIO


def extract_text_from_pptx(contents: bytes) -> str:
    presentation = Presentation(BytesIO(contents))

    slides_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())

        if slide_lines:
            slides_text.append(
                f"Slide {slide_number}:\n" + "\n".join(slide_lines)
            )

    return "\n\n".join(slides_text)
